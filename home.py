import pandas as pd
import streamlit as st
from io import BytesIO
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import matplotlib.pyplot as plt
from matplotlib import rcParams, font_manager as fm

# Load font dari file .otf
font_path = "VAGRoundedStd-Light.otf"
font_prop = fm.FontProperties(fname=font_path)

def apply_font_to_ax(ax, font):
    ax.title.set_fontproperties(font)
    ax.xaxis.label.set_fontproperties(font)
    ax.yaxis.label.set_fontproperties(font)
    for label in ax.get_xticklabels() + ax.get_yticklabels():
        label.set_fontproperties(font)

# Claim data functions
def filter_claim_data(df):
    return df[df['ClaimStatus'] == 'R']

def remove_duplicate_claims(df):
    dups = df[df.duplicated(subset='ClaimNo', keep=False)]
    if not dups.empty:
        st.write("Duplicated ClaimNo values:")
        st.dataframe(dups[['ClaimNo']].drop_duplicates())
    return df.drop_duplicates(subset='ClaimNo', keep='last')

def process_claim_data(df):
    df = filter_claim_data(df)
    df = remove_duplicate_claims(df)
    for col in ["TreatmentStart", "TreatmentFinish", "Date"]:
        df[col] = pd.to_datetime(df[col], errors='coerce')
        if df[col].isnull().any():
            st.warning(f"Invalid date values detected in column '{col}'. Coerced to NaT.")

    # Build standardized template:
    return pd.DataFrame({
        "No": range(1, len(df) + 1),
        "Policy No": df["PolicyNo"],
        "Client Name": df["ClientName"],
        "Claim No": df["ClaimNo"],
        "Member No": df["MemberNo"],
        "Emp ID": df["EmpID"],
        "Emp Name": df["EmpName"],
        "Patient Name": df["PatientName"],
        "Membership": df["Membership"],
        "Product Type": df["ProductType"],
        "Claim Type": df["ClaimType"],
        "Room Option": df["RoomOption"].fillna('').astype(str).str.upper().str.replace(r"\s+", "", regex=True),
        "Area": df["Area"],
        "Plan": df["PPlan"],
        "Diagnosis": df["PrimaryDiagnosis"].str.upper(),
        "Treatment Place": df["TreatmentPlace"].str.upper(),
        "Treatment Start": df["TreatmentStart"],
        "Treatment Finish": df["TreatmentFinish"],
        "Settled Date": df["Date"],
        "Year": df["Date"].dt.year,
        "Month": df["Date"].dt.month,
        "Length of Stay": df["LOS"],
        "Sum of Billed": df["Billed"],
        "Sum of Accepted": df["Accepted"],
        "Sum of Excess Coy": df["ExcessCoy"],
        "Sum of Excess Emp": df["ExcessEmp"],
        "Sum of Excess Total": df["ExcessTotal"],
        "Sum of Unpaid": df["Unpaid"]
    })

# Benefit data functions
def filter_benefit_data(df):
    if 'Status_Claim' in df.columns:
        return df[df['Status_Claim'] == 'R']
    elif 'Status Claim' in df.columns:
        return df[df['Status Claim'] == 'R']
    else:
        st.warning("Column 'Status Claim' not found. Data not filtered.")
        return df

def process_benefit_data(df):
    df = filter_benefit_data(df)
    df.columns = df.columns.str.strip()
    for col in df.columns:
        if df[col].dtype == "object":
            df[col] = df[col].astype(str).str.strip()
    rename_mapping = {
        'ClientName': 'Client Name',
        'PolicyNo': 'Policy No',
        'ClaimNo': 'Claim No',
        'MemberNo': 'Member No',
        'PatientName': 'Patient Name',
        'EmpID': 'Emp ID',
        'EmpName': 'Emp Name',
        'ClaimType': 'Claim Type',
        'TreatmentPlace': 'Treatment Place',
        'RoomOption': 'Room Option',
        'TreatmentRoomClass': 'Treatment Room Class',
        'TreatmentStart': 'Treatment Start',
        'TreatmentFinish': 'Treatment Finish',
        'ProductType': 'Product Type',
        'BenefitName': 'Benefit Name',
        'PaymentDate': 'Payment Date',
        'ExcessTotal': 'Excess Total',
        'ExcessCoy': 'Excess Coy',
        'ExcessEmp': 'Excess Emp'
    }
    df = df.rename(columns=rename_mapping)
    return df.drop(columns=["Status_Claim", "BAmount"], errors='ignore')

# Save to excel
def save_to_excel(claim_df, benefit_df, summary_top_df, claim_ratio_df, filename):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        workbook = writer.book
        workbook.formats[0].set_font_name('VAG Rounded Std Light')

        # Define excel formats:
        bold_border = workbook.add_format({'bold': True, 'border': 1, 'font_name': 'VAG Rounded Std Light'})
        plain_border = workbook.add_format({'border': 1, 'font_name': 'VAG Rounded Std Light'})
        header_border = workbook.add_format({'bold': True, 'border': 1, 'align': 'center', 'font_name': 'VAG Rounded Std Light'})

        # Summary sheet:
        summary_sheet = workbook.add_worksheet("Summary")
        summary_sheet.hide_gridlines(2)
        row = 0
        # Write summary statistics:
        for _, data in summary_top_df.iterrows():
            summary_sheet.write(row, 0, data["Metric"], bold_border)
            summary_sheet.write(row, 1, data["Value"], plain_border)
            row += 1
        # Insert blank row without borders (seperate sum stats & CR)
        summary_sheet.write(row, 0, "")
        summary_sheet.write(row, 1, "")
        row += 1

        # Write header for Claim Ratio table
        cr_columns = ["Company", "Net Premi", "Billed", "Unpaid", "Excess Total", "Excess Coy", "Excess Emp", "Claim", "CR", "Est Claim"]
        for col, header in enumerate(cr_columns):
            summary_sheet.write(row, col, header, header_border)
        row += 1

        # Write Claim Ratio data rows
        for _, data in claim_ratio_df.iterrows():
            for col, header in enumerate(cr_columns):
                summary_sheet.write(row, col, data.get(header, ""), plain_border)
            row += 1

        # Claim (SC) sheet
        claim_df.to_excel(writer, index=False, sheet_name='SC')
        ws_claim = writer.sheets["SC"]
        ws_claim.hide_gridlines(2)
        rows_claim, cols_claim = claim_df.shape[0] + 1, claim_df.shape[1]
        ws_claim.conditional_format(0, 0, rows_claim - 1, cols_claim - 1,
                                     {'type': 'no_errors', 'format': plain_border})
        for col_num, value in enumerate(claim_df.columns.values):
            ws_claim.write(0, col_num, value, header_border)

        # Benefit sheet
        benefit_df.to_excel(writer, index=False, sheet_name='Benefit')
        ws_benefit = writer.sheets["Benefit"]
        ws_benefit.hide_gridlines(2)
        rows_benefit, cols_benefit = benefit_df.shape[0] + 1, benefit_df.shape[1]
        ws_benefit.conditional_format(0, 0, rows_benefit - 1, cols_benefit - 1,
                                      {'type': 'no_errors', 'format': plain_border})
        for col_num, value in enumerate(benefit_df.columns.values):
            ws_benefit.write(0, col_num, value, header_border)


        writer.close()
    output.seek(0)
    return output, filename

def save_table_as_image(df, filename):
    fig, ax = plt.subplots(figsize=(len(df.columns) * 2.5, len(df) * 0.6 + 1), dpi=150)
    ax.axis('off')

    table = ax.table(
        cellText=df.values,
        colLabels=df.columns,
        cellLoc='center',
        loc='center'
    )

    # Styling
    for (i, j), cell in table.get_celld().items():
        cell.set_edgecolor('black')
        cell.set_linewidth(1)
        if i == 0:
            cell.set_facecolor('#0070C0')
            cell.get_text().set_color('white')
            cell.get_text().set_weight('bold')
        else:
            cell.set_facecolor('#fcfcfa' if i % 2 == 0 else 'white')
            cell.get_text().set_color('black')

    table.auto_set_font_size(False)
    table.set_fontsize(11)
    table.scale(1, 1.5)

    plt.tight_layout()
    fig.savefig(filename, bbox_inches='tight')
    plt.close(fig)

# Streamlit APP UI
st.title("Template - Standardisasi Report")

uploaded_claim = st.file_uploader("Upload Claim Data", type=["csv"], key="claim")
uploaded_claim_ratio = st.file_uploader("Upload Claim Ratio Data", type=["xlsx"], key="claim_ratio")
uploaded_benefit = st.file_uploader("Upload Benefit Data", type=["csv"], key="benefit")

if uploaded_claim and uploaded_claim_ratio and uploaded_benefit:
    # Process claim data
    raw_claim = pd.read_csv(uploaded_claim)
    st.write("Processing Claim Data...")
    claim_transformed = process_claim_data(raw_claim)
    st.subheader("Claim Data Preview:")
    st.dataframe(claim_transformed.head())

    # Process claim ratio data
    claim_ratio_raw = pd.read_excel(uploaded_claim_ratio)
    policy_list = claim_transformed["Policy No"].unique().tolist()
    claim_ratio_filtered = claim_ratio_raw[claim_ratio_raw["Policy No"].isin(policy_list)]
    claim_ratio_unique = claim_ratio_filtered.drop_duplicates(subset="Policy No")
    desired_cols = ['Company', 'Net Premi', 'Billed', 'Unpaid',
                    'Excess Total', 'Excess Coy', 'Excess Emp', 'Claim', 'CR', 'Est CR Total']
    missing_cols = [col for col in desired_cols if col not in claim_ratio_unique.columns]
    if missing_cols:
        st.warning(f"Missing columns in Claim Ratio Data: {missing_cols}")
    claim_ratio_unique = claim_ratio_unique[[col for col in desired_cols if col in claim_ratio_unique.columns]]
    claim_ratio_unique = claim_ratio_unique.rename(columns={'Est CR Total': 'Est Claim'})
    summary_cr_df = claim_ratio_unique[['Company', 'Net Premi', 'Billed', 'Unpaid',
                                         'Excess Total', 'Excess Coy', 'Excess Emp', 'Claim', 'CR', 'Est Claim']]
    st.subheader("Claim Ratio Data Preview (unique by Policy No):")
    st.dataframe(summary_cr_df.head())

    # Process benefit data
    raw_benefit = pd.read_csv(uploaded_benefit)
    st.write("Processing Benefit Data...")
    benefit_transformed = process_benefit_data(raw_benefit)
    claim_no_list = claim_transformed["Claim No"].unique().tolist()
    if "ClaimNo" in benefit_transformed.columns:
        benefit_transformed = benefit_transformed[benefit_transformed["ClaimNo"].isin(claim_no_list)]
    elif "Claim No" in benefit_transformed.columns:
        benefit_transformed = benefit_transformed[benefit_transformed["Claim No"].isin(claim_no_list)]
    else:
        st.warning("Column 'ClaimNo' not found in Benefit data; skipping filtering based on ClaimNo.")
    st.subheader("Benefit Data Preview:")
    st.dataframe(benefit_transformed.head())

    # Prepare Summary Top Section (Claim Stats + Overall Claim Ratio)
    total_claims   = len(claim_transformed)
    total_billed   = int(claim_transformed["Sum of Billed"].sum())
    total_accepted = int(claim_transformed["Sum of Accepted"].sum())
    total_excess   = int(claim_transformed["Sum of Excess Total"].sum())
    total_unpaid   = int(claim_transformed["Sum of Unpaid"].sum())
    claim_summary_data = {
        "Metric": ["Total Claims", "Total Billed", "Total Accepted", "Total Excess", "Total Unpaid"],
        "Value": [f"{total_claims:,}", f"{total_billed:,}", f"{total_accepted:,}",
                  f"{total_excess:,}", f"{total_unpaid:,}"]
    }
    claim_summary_df = pd.DataFrame(claim_summary_data)

    if "Claim" in claim_ratio_unique.columns and "Net Premi" in claim_ratio_unique.columns:
        total_claim_ratio_claim = claim_ratio_unique["Claim"].sum()
        total_net_premi = claim_ratio_unique["Net Premi"].sum()
        overall_cr = (total_claim_ratio_claim / total_net_premi) * 100 if total_net_premi != 0 else 0
        claim_ratio_overall = pd.DataFrame({"Metric": ["Claim Ratio (%)"],
                                            "Value": [f"{overall_cr:.2f}%"]})
    else:
        claim_ratio_overall = pd.DataFrame({"Metric": ["Claim Ratio (%)"], "Value": ["N/A"]})

    summary_top_df = pd.concat([claim_summary_df, claim_ratio_overall], ignore_index=True)

    # Download the Excel file
    st.subheader("Download Processed Data")
    filename_input = st.text_input("Enter the Excel file name (without extension):", "SC & Benefit - - YTD")
    if filename_input:
        excel_file, final_filename = save_to_excel(claim_transformed, benefit_transformed,
                                                   summary_top_df, summary_cr_df, filename_input + ".xlsx")
        st.download_button(
            label="Download Processed Excel File",
            data=excel_file,
            file_name=final_filename,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # Display Visualizations Section (Moved after download button)
    st.subheader("Visualizations")

    # Section 1: Summary Metrics Visualization
    st.subheader("Summary Metrics")
    # Prepare data for horizontal view
    metrics_row1 = summary_top_df.iloc[0:3]
    metrics_row2 = summary_top_df.iloc[3:]
    
    def display_metric(metric_name, metric_value):
        try:
            formatted_value = f"{int(float(str(metric_value).replace(',', ''))):,}"
        except:
            formatted_value = metric_value
        st.markdown(
            f"<p style='color: #0067B1; font-size: 18px; margin-bottom: 0;'>{metric_name}</p>"
            f"<p style='color: #0067B1; font-size: 24px; font-weight: bold; margin-top: 0;'>{metric_value}</p>",
            unsafe_allow_html=True
        )
    
    col1, col2, col3 = st.columns(3)
    with col1:
        display_metric(metrics_row1.iloc[0]['Metric'], metrics_row1.iloc[0]['Value'])
    with col2:
        display_metric(metrics_row1.iloc[1]['Metric'], metrics_row1.iloc[1]['Value'])
    with col3:
        display_metric(metrics_row1.iloc[2]['Metric'], metrics_row1.iloc[2]['Value'])
    
    col4, col5, col6 = st.columns(3)
    with col4:
        display_metric(metrics_row2.iloc[0]['Metric'], metrics_row2.iloc[0]['Value'])
    with col5:
        display_metric(metrics_row2.iloc[1]['Metric'], metrics_row2.iloc[1]['Value'])
    with col6:
        display_metric(metrics_row2.iloc[2]['Metric'], metrics_row2.iloc[2]['Value'])
    
    # Bersihkan nilai dan siapkan tabel
    summary_top_df['Formatted'] = summary_top_df['Value'].apply(
        lambda v: f"{int(float(str(v).replace(',', ''))):,}"
        if str(v).replace(',', '').replace('.', '').isdigit() else v
    )
    cell_text = summary_top_df[['Metric', 'Formatted']].values.tolist()
    
    # ← Ukuran figur diperbesar agar jelas di PPT
    fig, ax = plt.subplots(figsize=(14, 3), dpi=150)
    ax.axis('off')
    
    table = ax.table(
        cellText=cell_text,
        colLabels=['Metric', 'Value'],
        cellLoc='center',
        loc='center'
    )
    
    # Styling tabel & terapkan VAG Rounded Std Light
    table.auto_set_font_size(False)
    # Anda bisa atur ukuran header dan isi sel berbeda jika suka
    HEADER_FSIZE = 16
    CELL_FSIZE   = 14
    table.scale(1.5, 2)
    
    for (i, j), cell in table.get_celld().items():
        cell.set_edgecolor('black')
        cell.set_linewidth(1)
        txt = cell.get_text()
        # terapkan font_prop untuk setiap sel
        txt.set_fontproperties(font_prop)
        if i == 0:
            cell.set_facecolor('#0070C0')
            txt.set_color('white')
            txt.set_weight('bold')
            txt.set_fontsize(HEADER_FSIZE)
        else:
            cell.set_facecolor('#fcfcfa' if i % 2 == 0 else 'white')
            txt.set_color('black')
            txt.set_fontsize(CELL_FSIZE)
    
    # Simpan sebagai PNG
    summary_path = "section1_summary_metrics.png"
    fig.savefig(summary_path, bbox_inches='tight')
    plt.close(fig)
    
    
    # Validasi
    import os
    if os.path.exists(summary_path):
        st.success(f"Summary metrics berhasil disimpan sebagai gambar: `{summary_path}`")
    else:
        st.error("Gagal menyimpan summary metrics.")

    st.image(
    summary_path,
    caption="Summary Metrics Table (akan dimasukkan ke PPT)",
    use_container_width=True
    )

        
    # Claim Ratio Summary Table (Enhanced Display & PNG output)
    st.subheader("Claim Ratio Summary Table")
    
    def save_claim_ratio_table_image(df, filename):
        headers = df.columns.tolist()
        cell_text = []
    
        for _, row in df.iterrows():
            formatted = []
            for col in headers:
                v = row[col]
                # Coba convert ke float bila perlu
                try:
                    num = float(v)
                except:
                    num = None
    
                if col == 'CR' and num is not None:
                    formatted.append(f"{num:.2f}%")
                elif col == 'Est Claim' and num is not None:
                    formatted.append(f"{num:,.2f}")
                elif num is not None:
                    formatted.append(f"{int(num):,}")
                else:
                    # fallback ke string as-is
                    formatted.append(str(v))
            cell_text.append(formatted)
    
        # ukuran figure dinamis
        ncols = len(headers)
        nrows = len(cell_text) + 1

        fig_width = max(16, ncols * 2.5)
        fig_height = max(4, nrows * 0.6)
        fig, ax = plt.subplots(figsize=(fig_width, fig_height), dpi=150)
        ax.axis('off')

        tbl = ax.table(
            cellText=cell_text,
            colLabels=headers,
            cellLoc='center',
            loc='center',
            colWidths=[0.4] + [0.6/(ncols-1)]*(ncols-1)
        )
        tbl.auto_set_column_width(col=list(range(ncols)))
        tbl.auto_set_font_size(False)
        # Setelah semuanya siap, set height tiap cell
        # row_height dalam axis fraction: tinggi axes dikali frac = cell height
        axis_pos = ax.get_position()  # BBox: [xmin, ymin, xmax, ymax]
        axis_h = axis_pos.height
        # kita bagi axis height ke nrows, lalu kurangi sedikit agar ada jarak
        row_height = axis_h / nrows * 0.9
        HEADER_FS, CELL_FS = 18, 16
    
        for (i, j), cell in tbl.get_celld().items():
            txt = cell.get_text()
            txt.set_fontproperties(font_prop)
            if i == 0:
                cell.set_facecolor('#0070C0')
                txt.set_color('white')
                txt.set_weight('bold')
                txt.set_fontsize(HEADER_FS)
            else:
                cell.set_facecolor('#fcfcfa' if i % 2 == 0 else 'white')
                txt.set_color('black')
                txt.set_fontsize(CELL_FS)
            cell.set_edgecolor('black')
            cell.set_linewidth(1)
            cell.set_height(row_height)
    
        plt.tight_layout()
        fig.savefig(filename, bbox_inches='tight')
        plt.close(fig)
        
    # Simpan dan tampilkan PNG
    summary_table_name = "claim_ratio_table.png"
    save_claim_ratio_table_image(summary_cr_df, summary_table_name)
    if os.path.exists(summary_table_name):
        st.success(f"Tabel berhasil disimpan sebagai gambar: `{summary_table_name}`")
        st.image(summary_table_name, caption="Claim Ratio Summary Table", use_container_width=True)
    else:
        st.error("Gagal menyimpan tabel sebagai gambar.")



    
    # Section 2: Claim per Membership (Pie Chart)
    st.subheader("Claim Count per Membership Type")
    pie_path = None
    
    if 'Membership' in claim_transformed.columns:
        # Hitung jumlah klaim per Membership
        mc = claim_transformed['Membership'].value_counts()
        labels = mc.index.tolist()
        sizes = mc.values.tolist()

        label_map = {
            "1. EMP": "Employee",
            "2. SPO": "Spouse",
            "3. CHI": "Children"
        }
        mapped_labels = [label_map.get(l, l) for l in labels]
    
        # Format fungsi autopct: tampilkan persen + jumlah klaim
        def format_autopct(pct, allvals):
            absolute = int(round(pct/100.*sum(allvals)))
            return f"{pct:.1f}%\n({absolute:,})"
    
        # Buat figure lebih besar dan simetris
        fig, ax = plt.subplots(figsize=(7, 5))
        apply_font_to_ax(ax, font_prop)
        
        wedges, texts, autotexts = ax.pie(
            sizes,
            labels=None,
            colors=['#1f77b4', '#4e91c7', '#a6c8ea'],
            autopct=lambda pct: format_autopct(pct, sizes),
            textprops=dict(color="black", fontsize=10, fontproperties=font_prop),
            startangle=90,
        )
    
        # Tambahkan legend di kanan pie chart
        legend = ax.legend(wedges, mapped_labels, title="Membership", loc="center left", bbox_to_anchor=(1, 0.5), title_fontproperties=font_prop)

        # Terapkan font ke legend texts & title
        for txt in legend.get_texts():
            txt.set_fontproperties(font_prop)
        legend.get_title().set_fontproperties(font_prop)
        
        ax.axis('equal')  # Lingkaran bulat, bukan elips
    
        pie_path = "section2_membership.png"
        fig.savefig(pie_path, bbox_inches='tight')
        st.pyplot(fig)
        plt.close(fig)
    
        # Cek apakah file berhasil disimpan
        if os.path.exists(pie_path):
            st.success(f"Chart berhasil disimpan sebagai gambar: `{pie_path}`")
        else:
            st.error("Gagal menyimpan chart sebagai gambar.")
    else:
        st.warning("'Membership' column not found")
    
    # Section 3: Claim Count per Plan
    st.subheader("Claim Count per Plan")
    bar_path = None
    if 'Plan' in claim_transformed.columns:
        # Hitung jumlah klaim per Plan
        pc = claim_transformed['Plan'].value_counts().sort_index()
        plans = pc.index.tolist()
        counts = pc.values.tolist()
    
        # Buat bar chart
        fig3, ax3 = plt.subplots(figsize=(10, 7))
        apply_font_to_ax(ax3, font_prop)  # set ke axis default
        bars = ax3.bar(plans, counts, color='#1f77b4')
    
        # Tampilkan count sebagai label bar
        ax3.bar_label(bars, labels=[f"{c:,}" for c in counts], padding=3, color='black', fontproperties=font_prop)
    
        ax3.set_ylabel("Number of Claims", color='black', fontproperties=font_prop)
        plt.xticks(rotation=45, ha='right', fontproperties=font_prop)
    
        bar_path = "section3_plan.png"
        fig3.savefig(bar_path, bbox_inches='tight')
        st.pyplot(fig3)
        plt.close(fig3)
        if os.path.exists(bar_path):
            st.success(f"Tabel berhasil disimpan sebagai gambar: `{bar_path}`")
        else:
            st.error("Gagal menyimpan tabel sebagai gambar.")

    else:
        st.warning("'Plan' column not found")


    # ─── Section 4: Claim Billed by Month and Product Type ────────────────────────
    st.subheader("Claim Billed by Month and Product Type")
    month_prod_path = None
    
    if 'Settled Date' in claim_transformed.columns and 'Product Type' in claim_transformed.columns:
        # … proses pivot seperti biasa …
         # 1) Buat kolom Settled Month
        claim_transformed['Settled Month'] = claim_transformed['Settled Date'].dt.strftime("%b'%y")

        mbp = (
            claim_transformed
            .groupby(['Settled Month', 'Product Type'])['Sum of Billed']
            .sum()
            .reset_index()
        )
        # 3) Tentukan order bulan
        order = (
            claim_transformed['Settled Date']
            .dt.to_period('M')
            .sort_values()
            .dt.strftime("%b'%y")
            .unique()
        )
        
        # Lalu kamu urutkan dan pivot:
        mbp['Settled Month'] = pd.Categorical(mbp['Settled Month'], categories=order, ordered=True)
        mbp = mbp.sort_values('Settled Month')
        pivot = mbp.pivot(index='Settled Month', columns='Product Type', values='Sum of Billed').fillna(0)
    
        # Buat bar chart dan dapatkan ax
        fig, ax = plt.subplots(figsize=(12, 8))
        pivot.plot(kind='bar', ax=ax)
    
        # Terapkan font ke semua elemen axis
        apply_font_to_ax(ax, font_prop)

        # Axis labels
        ax.set_xlabel(
            "Settled Month",
            fontproperties=font_prop,
            fontsize=18,
            labelpad=15
        )
        ax.set_ylabel(
            "Sum of Billed",
            fontproperties=font_prop,
            fontsize=18,
            labelpad=15
        )
    
        # Format tick labels
        ax.tick_params(
            axis='x',
            rotation=45,
            labelsize=20,
            pad=10
        )
        # for lbl in ax.get_xticklabels():
        #     lbl.set_fontsize(20)
        #     lbl.set_fontproperties(font_prop)

        ax.tick_params(
            axis='y',
            labelsize=20,
            pad=10
        )
        # for lbl in ax.get_yticklabels():
        #     lbl.set_fontsize(20)
        #     lbl.set_fontproperties(font_prop)

        plt.setp(ax.get_xticklabels(), fontsize=20, fontproperties=font_prop)
        plt.setp(ax.get_yticklabels(), fontsize=20, fontproperties=font_prop)

            
        # Jika ada legend, terapkan font juga
        legend = ax.get_legend()
        if legend:
            legend.set_title("Product Type", prop=font_prop)
            for text in legend.get_texts():
                text.set_fontproperties(font_prop)
                text.set_fontsize(18)
        plt.tight_layout() 
        
        # Simpan & tampilkan
        month_prod_path = "section4_month_product.png"
        fig.savefig(month_prod_path, bbox_inches='tight')
        st.pyplot(fig)
        plt.close(fig)
    
        # Validasi simpan
        if os.path.exists(month_prod_path):
            st.success(f"Tabel berhasil disimpan sebagai gambar: `{month_prod_path}`")
        else:
            st.error("Gagal menyimpan tabel sebagai gambar.")
    else:
        st.warning("'Settled Date' or 'Product Type' column not found")


    
   # Tabel detail
    
    # Format angka dengan koma ribuan dan tanpa index tambahan
    pivot_formatted = pivot.copy()
    pivot_formatted = pivot_formatted.applymap(lambda x: f"{int(x):,}" if not pd.isna(x) else "")
    
    # Gabungkan kembali Settled Month ke dalam dataframe tanpa index
    final_table = pivot_formatted.reset_index()
    final_table = final_table.drop(columns=['index'], errors='ignore')
    
    # Tampilkan tanpa index tambahan
    st.dataframe(final_table, use_container_width=True, hide_index=True)
    table_filename = "section4_month_product_table.png"
    save_table_as_image(final_table, table_filename)
    if os.path.exists(table_filename):
        st.success(f"Tabel berhasil disimpan sebagai gambar: `{table_filename}`")
    else:
        st.error("Gagal menyimpan tabel sebagai gambar.")

    
   # Section 5: Top 10 Diagnoses by Product Type
    st.subheader("Top 10 Diagnoses by Product Type")
    diag_path = []
    
    for product in claim_transformed['Product Type'].unique():
        dfp = (
            claim_transformed[claim_transformed['Product Type'] == product]
            .groupby('Diagnosis')['Sum of Billed']
            .agg(['sum', 'count'])
            .rename(columns={'sum': 'Amount', 'count': 'Qty'})
            .reset_index()
        )
        dfp['Amount'] /= 1_000_000  # dalam juta
        top10 = dfp.sort_values('Amount', ascending=False).head(10).iloc[::-1]
    
        n = len(top10)
        fig, ax = plt.subplots(figsize=(15, 0.65 * n + 2))
        apply_font_to_ax(ax, font_prop)  # set ke axis default

        # Dynamic font size
        max_label_length = max(top10['Diagnosis'].str.len())

        # Tetapkan font minimal yang besar dan skala naik sesuai panjang label
        if   max_label_length > 50:
            label_font = 24
        elif max_label_length > 30:
            label_font = 26
        else:
            label_font = 28
        
        # Value font selalu 2 poin lebih kecil dari label_font
        value_font = label_font - 2
        bar_height = 0.35
    
        y = range(n)
    
        # Bar Amount (di atas)
        ax.barh([i + bar_height/2 for i in y], top10['Amount'], height=bar_height,
                color='#1f77b4', label='Amount (mil)', alpha=0.9)
    
        # Bar Qty (di bawah)
        ax.barh([i - bar_height/2 for i in y], top10['Qty'], height=bar_height,
                color='#a6c8ea', label='Qty', alpha=0.9)
    
        # Label angka di samping bar
        for i, (amt, qty) in enumerate(zip(top10['Amount'], top10['Qty'])):
            ax.text(amt + 0.5, i + bar_height/2, f'{amt:,.0f}', va='center', fontsize=value_font, color='black', fontproperties=font_prop)
            ax.text(qty + 0.5, i - bar_height/2, f'{qty:,}', va='center', fontsize=value_font, color='black', fontproperties=font_prop)
    
        # Label sumbu Y
        ax.set_yticks(y)
        ax.set_yticklabels(top10['Diagnosis'], fontsize=label_font, fontproperties=font_prop)
    
        # Judul dan axis
        ax.set_title(f"{product}", fontsize=label_font + 4, weight='bold', fontproperties=font_prop)
        ax.set_xlabel("Value", fontsize=label_font, fontproperties=font_prop)
        ax.set_ylabel("Value", fontsize=label_font, fontproperties=font_prop)
        ax.tick_params(axis='x', labelsize=20)
        ax.tick_params(axis='y', labelsize=20)
    
        # Legend
        ax.legend(loc='lower right', fontsize=label_font, frameon=True)
    
        # Otomatis set lebar sumbu X
        max_val = max(top10['Amount'].max(), top10['Qty'].max())
        ax.set_xlim(0, max_val * 1.3)  # tambahkan 30% untuk label
    
        plt.tight_layout(pad=2)
    
        path = f"section5_diag_{product}.png"
        fig.savefig(path, bbox_inches='tight')
        st.pyplot(fig)
        plt.close(fig)
        if os.path.exists(path):
            st.success(f"Tabel berhasil disimpan sebagai gambar: `{path}`")
        else:
            st.error("Gagal menyimpan tabel sebagai gambar.")
    
        diag_path.append((product, path))


   # ─── Section 6: Top 10 Treatment Places by Claim Type ────────────────────────
    st.subheader("Top 10 Treatment Places by Claim Type")
    tp_path = []
    
    for claim_type in claim_transformed['Claim Type'].unique():
        dfp = (
            claim_transformed[claim_transformed['Claim Type'] == claim_type]
            .groupby('Treatment Place')['Sum of Billed']
            .agg(['sum', 'count'])
            .rename(columns={'sum': 'Amount', 'count': 'Qty'})
            .reset_index()
        )
        dfp['Amount'] /= 1_000_000  # dalam juta
        top10 = dfp.sort_values('Amount', ascending=False).head(10).iloc[::-1]
    
        n = len(top10)
        fig, ax = plt.subplots(figsize=(15, 0.65 * n + 2))
    
        # Font dan tinggi bar fleksibel
        max_label_length = max(top10['Treatment Place'].str.len())
        label_font = 14 if max_label_length > 40 else 15 if max_label_length > 30 else 16
        value_font = max(12, label_font - 2)
        bar_height = 0.35
        y = range(n)
    
        # Bar Amount (di atas)
        ax.barh([i + bar_height/2 for i in y], top10['Amount'], height=bar_height,
                color='#1f77b4', label='Amount (mil)', alpha=0.9)
    
        # Bar Qty (di bawah)
        ax.barh([i - bar_height/2 for i in y], top10['Qty'], height=bar_height,
                color='#a6c8ea', label='Qty', alpha=0.9)
    
        # Label angka
        for i, (amt, qty) in enumerate(zip(top10['Amount'], top10['Qty'])):
            ax.text(amt + 0.5, i + bar_height/2, f'{amt:,.1f}', va='center', fontsize=value_font, color='black')
            ax.text(qty + 0.5, i - bar_height/2, f'{qty:,}', va='center', fontsize=value_font, color='black')
    
        # Label tempat perawatan
        ax.set_yticks(y)
        ax.set_yticklabels(top10['Treatment Place'], fontsize=label_font)
    
        # Judul, axis, dan legend
        ax.set_title(f"Top 10 Treatment Places: {claim_type}", fontsize=label_font + 4, weight='bold')
        ax.set_xlabel("Value", fontsize=label_font)
        ax.tick_params(axis='x', labelsize=label_font)
        ax.legend(loc='lower right', fontsize=label_font, frameon=True)
    
        # Batas x agar semua label muat
        max_val = max(top10['Amount'].max(), top10['Qty'].max())
        ax.set_xlim(0, max_val * 1.3)
    
        plt.tight_layout(pad=2)
    
        path = f"section6_tp_{claim_type}.png"
        fig.savefig(path, bbox_inches='tight')
        st.pyplot(fig)
        plt.close(fig)
        if os.path.exists(path):
            st.success(f"Tabel berhasil disimpan sebagai gambar: `{path}`")
        else:
            st.error("Gagal menyimpan tabel sebagai gambar.")
    
        tp_path.append((claim_type, path))


    # Section 7: Top 10 Employee
    st.subheader("Top 10 Employees by Number of Claims")

    df_emp = claim_transformed.copy()
    
    # Group and summarize
    top_10_emp_summary = (
        df_emp.groupby(['Emp Name', 'Plan'])
              .agg(
                  Total_Claims=('Emp Name', 'count'),
                  Total_Billed=('Sum of Billed', 'sum')
              )
              .reset_index()
    )
    
    # Show warning if empty
    if top_10_emp_summary.empty:
        st.warning("No employee data available.")
    else:
        # Sort & take top 10
        top_10_emp_summary = top_10_emp_summary.sort_values(by='Total_Claims', ascending=False).head(10)
    
        # Rename and reorder
        top_10_emp_summary = top_10_emp_summary.rename(columns={
            'Emp Name': 'Employee',
            'Total_Claims': 'Total Claims',
            'Total_Billed': 'Total Billed'
        })[['Employee', 'Plan', 'Total Claims', 'Total Billed']]
    
        # Format numbers
        top_10_emp_summary['Total Claims'] = top_10_emp_summary['Total Claims'].map('{:,}'.format)
        top_10_emp_summary['Total Billed'] = top_10_emp_summary['Total Billed'].map('{:,}'.format)
    
        # Styled HTML table with thick visible borders
        def render_styled_table(df):
            html = """
            <style>
            table {
                border-collapse: collapse;
                width: 100%;
                font-family: Arial, sans-serif;
            }
            tr:nth-child(even) {
                background-color: #f5f5f5;
            }
            </style>
            <table>
                <thead>
                    <tr>
            """
            # Tambahkan header dengan inline border
            for col in df.columns:
                html += f"<th style='border: 1px solid #333; background-color: #0067B1; color: white; padding: 8px;'>{col}</th>"
            html += "</tr></thead><tbody>"
        
            # Tambahkan isi tabel dengan border per cell
            for _, row in df.iterrows():
                html += "<tr>"
                for item in row:
                    html += f"<td style='border: 1px solid #333; color: black; padding: 8px; text-align: center;'>{item}</td>"
                html += "</tr>"
        
            html += "</tbody></table>"
            return html
    
        # Render in Streamlit
        st.markdown(render_styled_table(top_10_emp_summary), unsafe_allow_html=True)
        # Simpan sebagai gambar
        table_path = "section7_top10_employees.png"
        save_table_as_image(top_10_emp_summary, table_path)
        
        # Beri info apakah berhasil disimpan
        if os.path.exists(table_path):
            st.success(f"Tabel berhasil disimpan sebagai gambar: `{table_path}`")
        else:
            st.error("Gagal menyimpan tabel sebagai gambar.")


    # Gunakan layout 'Content Slide 4' (index ke-5)
    CONTENT_LAYOUT_INDEX = 5
    
    def add_title(slide, title_text):
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        else:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
            title_tf = title_shape.text_frame
            title_tf.text = title_text
            title_tf.paragraphs[0].font.size = Pt(24)
            title_tf.paragraphs[0].font.bold = True
    
    def create_ppt(path):
        prs = Presentation("tes.pptx")
        content_layout = prs.slide_layouts[5]  # Content Slide 4
    
        # Simpan slide baru ke list dulu
        new_slides = []
    
        def generate_slide(title, image_path):
            if os.path.exists(image_path):
                slide = prs.slides.add_slide(content_layout)
                add_title(slide, title)
                slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(7))
                new_slides.append(slide)
    
        # Kumpulkan semua slide baru
        generate_slide("Summary Metrics", "section1_summary_metrics.png")
        generate_slide("Claim Ratio Summary Table", "claim_ratio_table.png")
        generate_slide("Claim Count per Membership Type", "section2_membership.png")
        generate_slide("Claim Count per Plan", "section3_plan.png")
        generate_slide("Claim Billed by Month and Product Type", "section4_month_product.png")
        generate_slide("Claim Table by Month and Product Type", "section4_month_product_table.png")
        generate_slide("Top 10 Employees by Number of Claims", "section7_top10_employees.png")
    
        # Diagnosa dan Treatment dari loop
        for product, path_img in diag_path:
            if os.path.exists(path_img):
                slide = prs.slides.add_slide(content_layout)
                add_title(slide, f"Top 10 Diagnoses: {product}")
                slide.shapes.add_picture(path_img, Inches(1), Inches(1.5), width=Inches(7))
                new_slides.append(slide)
    
        for claim_type, path_img in tp_path:
            if os.path.exists(path_img):
                slide = prs.slides.add_slide(content_layout)
                add_title(slide, f"Top 10 Treatment Places: {claim_type}")
                slide.shapes.add_picture(path_img, Inches(1), Inches(1.5), width=Inches(7))
                new_slides.append(slide)
    
        for _ in range(len(new_slides)):
            slide = prs.slides[-1]
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            xml_slides.insert(1, slides[-1])  # << selalu di posisi ke-1

        prs.save(path)

    
    # =============================
    #       STREAMLIT UI
    # =============================
    st.markdown("---")
    st.subheader("📊 Generate PowerPoint Report")
    
    # ⬇️ Input untuk nama file
    ppt_filename_input = st.text_input("Enter PPT file name (without .pptx):", "Claim_Report")
    ppt_filename = (ppt_filename_input.strip() or "Claim_Report") + ".pptx"
    ppt_filepath = os.path.join(".", ppt_filename)

    # layout_names = [f"{i}: {layout.name}" for i, layout in enumerate(prs.slide_layouts)]
    # st.write("Available Layouts in Template:")
    # st.code("\n".join(layout_names))
    
    # ⬇️ Tombol untuk generate
    if st.button("Generate PPT"):
        create_ppt(ppt_filepath)
    
        if os.path.exists(ppt_filepath):
            with open(ppt_filepath, "rb") as f:
                st.success("✅ PPT generated successfully!")
                st.download_button(
                    label="⬇️ Download PPT",
                    data=f,
                    file_name=ppt_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
        else:
            st.error("❌ Failed to generate PPT.")
